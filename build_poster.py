"""Build the A0 portrait project poster (Climate Capacitor)."""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

REPO = r"D:/Climate-Modeling-Research-Project"

# ---- palette (hot/cold capacitor theme) ----
NAVY   = RGBColor(0x12, 0x23, 0x3A)
HEAT   = RGBColor(0xC1, 0x27, 0x2D)   # red  (positive charge / heat)
COLD   = RGBColor(0x1B, 0x6C, 0xA8)   # blue (negative charge / cold)
CARD   = RGBColor(0xF3, 0xF5, 0xF8)   # very light card
INK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTE   = RGBColor(0x5A, 0x66, 0x72)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)

W, H = 84.1, 118.9
prs = Presentation()
prs.slide_width  = Cm(W)
prs.slide_height = Cm(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes


def rect(l, t, w, h, fill, line=None, rounded=False):
    shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Cm(l), Cm(t), Cm(w), Cm(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,bold,color,italic)."""
    tb = shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, Cm(0.15))
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space
        p.space_after = Pt(4)
        for (txt, size, bold, color, *rest) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.name = "Arial"; r.font.color.rgb = color
            if rest and rest[0]:
                r.font.italic = True
    return tb


def img(path, l, t, w):
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
        h = w * ih / iw
    except Exception:
        h = w * 0.5
    shapes.add_picture(path, Cm(l), Cm(t), Cm(w), Cm(h))
    return h


# ============ HEADER ============
rect(0, 0, W, 11.5, NAVY)
# small + / - charge motif
c1 = shapes.add_shape(MSO_SHAPE.OVAL, Cm(75.5), Cm(1.3), Cm(3.2), Cm(3.2))
c1.fill.solid(); c1.fill.fore_color.rgb = HEAT; c1.line.fill.background(); c1.shadow.inherit=False
c2 = shapes.add_shape(MSO_SHAPE.OVAL, Cm(75.5), Cm(5.2), Cm(3.2), Cm(3.2))
c2.fill.solid(); c2.fill.fore_color.rgb = COLD; c2.line.fill.background(); c2.shadow.inherit=False
text(75.3, 1.2, 3.6, 3.4, [[("+", 40, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
text(75.3, 5.1, 3.6, 3.4, [[("−", 40, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

text(2, 1.2, 72, 5.5, [[("THE CLIMATE CAPACITOR", 74, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(2, 6.7, 72, 2.4, [[("Can extreme-weather disasters be predicted as “electrical breakdowns” of accumulated heat?",
                          30, False, ICE, True)]])
text(2, 9.1, 80, 2.0, [[("Project 26-1-R-14      ·      Ron Salama  &  Lior Dagash      ·      Supervisor: Dr. Zakharia Frenkel",
                          22, True, WHITE)]])

# ============ 1. BACKGROUND & ANALOGY ============
y = 12.3
rect(1.5, y, 81, 15.2, CARD, rounded=True)
text(2.3, y+0.4, 79, 1.6, [[("1 · BACKGROUND & THE CAPACITOR ANALOGY", 30, True, NAVY)]])
text(2.3, y+2.2, 79, 12.5, [
 [("Extreme weather — floods, storms, heatwaves — strikes where energy builds up and suddenly releases. "
   "Conventional forecasting simulates atmospheric fluid dynamics; this project tests a different, structural idea: "
   "the ", 20, False, INK), ("Climate Capacitor", 20, True, HEAT),
  (" analogy, which treats Earth’s surface as one giant capacitor.", 20, False, INK)],
 [("●  Thermal charge — ", 20, True, HEAT), ("temperature anomalies (hotter/colder than normal) accumulate in each grid cell.", 20, False, INK)],
 [("●  Dielectric (terrain) — ", 20, True, COLD), ("mountains and valleys resist the release of energy (permittivity ε).", 20, False, INK)],
 [("●  Breakdown — ", 20, True, NAVY), ("a disaster is the “discharge” where the gradient between adjacent hot and cold regions grows too steep for the terrain to contain.", 20, False, INK)],
 [("Core question: does this electrostatic lens reveal where and when disasters actually strike?", 20, False, MUTE, True)],
], space=1.05)

# ============ 2. HOW IT WORKS (pipeline) ============
y = 28.5
rect(1.5, y, 81, 13.0, NAVY, rounded=True)
text(2.3, y+0.4, 55, 1.6, [[("2 · HOW IT WORKS", 30, True, WHITE)]])
text(45, y+0.5, 37, 1.8, [[("breakdown field   E = ‖∇Q‖ / ε", 26, True, ICE)]], PP_ALIGN.RIGHT)
stages = ["Temperature\nanomaly", "Accumulate\ncharge  Q", "Terrain\npermittivity ε",
          "Breakdown\nfield  E", "Cluster into\nevents", "Validate vs\ndisasters"]
n = len(stages); gap = 0.6; bw = (81 - 2*0.8 - (n-1)*gap) / n; bx = 2.3; by = y+3.0; bh = 6.2
for i, s in enumerate(stages):
    col = HEAT if i in (0,1) else COLD if i==2 else NAVY
    ch = shapes.add_shape(MSO_SHAPE.CHEVRON, Cm(bx+i*(bw+gap)), Cm(by), Cm(bw+0.9), Cm(bh))
    ch.fill.solid(); ch.fill.fore_color.rgb = WHITE if i not in (0,1,2) else col
    ch.line.color.rgb = WHITE; ch.line.width = Pt(1.5); ch.shadow.inherit=False
    tf = ch.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    for j, ln in enumerate(s.split("\n")):
        pp = p if j==0 else tf.add_paragraph(); pp.alignment=PP_ALIGN.CENTER
        r=pp.add_run(); r.text=ln; r.font.size=Pt(15); r.font.bold=True; r.font.name="Arial"
        r.font.color.rgb = WHITE if i in (0,1,2) else NAVY
text(2.3, y+9.6, 79, 2.2, [[("Charge is accumulated over a sliding 30-day window; the breakdown field is its spatial "
                              "gradient divided by terrain permittivity. High-stress cells are clustered into events "
                              "and matched against real disasters.", 16, False, ICE)]])

# ============ 3 & 4 two columns ============
y = 43.0
colw = 39.75
rect(1.5, y, colw, 19.0, CARD, rounded=True)
rect(1.5+colw+1.0, y, colw, 19.0, CARD, rounded=True)
text(2.3, y+0.4, colw-1, 1.6, [[("3 · SYSTEM REQUIREMENTS", 28, True, NAVY)]])
text(2.3, y+2.1, colw-1.4, 16.5, [
 [("Functional", 20, True, HEAT)],
 [("●  Ingest 10 yrs of global ERA5 temperature + terrain", 17, False, INK)],
 [("●  Compute anomalies vs a day-of-year climatology", 17, False, INK)],
 [("●  Accumulate thermal charge; map terrain → ε", 17, False, INK)],
 [("●  Build the breakdown field; flag & cluster zones", 17, False, INK)],
 [("●  Cross-reference zones with a disaster database", 17, False, INK)],
 [("●  Generate global risk / validation maps", 17, False, INK)],
 [("Non-functional", 20, True, COLD)],
 [("●  Scale to global 10-yr data with no memory overflow", 17, False, INK)],
 [("●  A full run completes in minutes on normal hardware", 17, False, INK)],
 [("●  Modular & config-driven — any formula is swappable", 17, False, INK)],
 [("●  Auditable — trace any breakdown back to raw data", 17, False, INK)],
], space=1.05)

lx = 1.5+colw+1.0
text(lx+0.8, y+0.4, colw-1, 1.6, [[("4 · TOOLS & TECHNOLOGIES", 28, True, NAVY)]])
text(lx+0.8, y+2.1, colw-1.4, 16.5, [
 [("●  Language:  ", 18, True, INK), ("Python", 18, False, INK)],
 [("●  Data engine:  ", 18, True, INK), ("xarray + dask (chunked, lazy)", 18, False, INK)],
 [("●  Climate + terrain:  ", 18, True, INK), ("ERA5 via WeatherBench2 — public cloud, streamed (no bulk download)", 18, False, INK)],
 [("●  Disasters:  ", 18, True, INK), ("EM-DAT (dates) + GDIS (precise coordinates)", 18, False, INK)],
 [("●  Analysis:  ", 18, True, INK), ("SciPy, scikit-learn / ndimage (clustering), NumPy", 18, False, INK)],
 [("●  Stats & viz:  ", 18, True, INK), ("pandas, Matplotlib", 18, False, INK)],
 [("Six-stage pipeline, one config file", 18, True, COLD)],
 [("data → anomaly → charge → permittivity → breakdown → events → validation. Every "
   "knob (terrain weighting, timescale, thresholds) lives in one config, enabling systematic experiments.",
   16, False, MUTE)],
], space=1.05)

# ============ 5. RESULTS & METRICS ============
y = 63.5
rect(1.5, y, 81, 29.0, CARD, rounded=True)
text(2.3, y+0.4, 79, 1.6, [[("5 · RESULTS & METRICS", 30, True, NAVY)]])
# metric callouts
def callout(l, t, big, bigcol, label):
    text(l, t, 18.0, 3.0, [[(big, 46, True, bigcol)]], PP_ALIGN.CENTER)
    labs = [[(ln, 15, False, MUTE)] for ln in label.split("\n")]
    text(l, t+3.0, 18.0, 2.6, labs, PP_ALIGN.CENTER, space=1.0)
ry = y+2.3
callout(2.5, ry, "~1–3%", HEAT, "Recall achieved\n(target > 30%)")
callout(21.5, ry, "< 2%", HEAT, "Precision achieved\n(target > 10%)")
text(2.5, ry+6.0, 38, 8.0, [
 [("Statistically “significant” but practically negligible — huge p-values are a "
   "sample-size effect, not real skill. The analogy did ", 16, False, INK),
  ("not", 16, True, HEAT), (" meet any success target.", 16, False, INK)],
 [("Terrain (the “dielectric”) added almost nothing; and finer resolution, precise "
   "GDIS coordinates, and precursor-timing did not help either. The high-stress "
   "zones and real disasters occupy different places.", 16, False, MUTE)],
], space=1.05)
# figure: predicted vs actual
ph = img(os.path.join(REPO, "outputs/phase4/predicted_vs_actual.png"), 42.0, ry, 40.0)
text(42.0, ry+ph+0.1, 40, 1.8, [[("Predicted breakdown zones (red) vs. real disasters (blue): they mostly "
                                   "occupy different places — the core of the negative result.", 13, False, MUTE, True)]])

# ============ 6. CHALLENGES ============
y = 93.5
rect(1.5, y, 81, 15.0, NAVY, rounded=True)
text(2.3, y+0.4, 79, 1.6, [[("6 · CHALLENGES & CONSTRAINTS", 30, True, WHITE)]])
text(2.3, y+2.2, 79, 12.0, [
 [("●  Data volume (terabytes) — ", 18, True, ICE), ("streamed coarse, daily, single-variable from the cloud (~1 GB); no local hoarding.", 18, False, WHITE)],
 [("●  Runtime — ", 18, True, ICE), ("parallel cloud fetch + connected-component clustering → minutes, not hours.", 18, False, WHITE)],
 [("●  Disaster geolocation gap — ", 18, True, ICE), ("EM-DAT coords were mostly missing; joined GDIS coordinates to EM-DAT dates.", 18, False, WHITE)],
 [("●  Memory — ", 18, True, ICE), ("float32 + local caching; the fine 0.7° grid needs ≥ 32 GB RAM (it blue-screened a 16 GB laptop).", 18, False, WHITE)],
 [("●  Method artifacts — ", 18, True, ICE), ("km-based gradient + polar masking; the strict 100 km metric is near-impossible below grid scale.", 18, False, WHITE)],
], space=1.1)

# ============ CONCLUSION ============
y = 109.2
rect(0, y, W, H-y, HEAT)
text(2.3, y+0.5, 79, 1.6, [[("CONCLUSION", 28, True, WHITE)]])
text(2.3, y+2.1, 79, 6.5, [
 [("The Climate Capacitor analogy does not usefully predict real disasters — its high-stress zones and "
   "actual events largely diverge, and the terrain component adds little. This is a clear, ", 19, False, WHITE),
  ("valuable negative result", 19, True, WHITE),
  (": it maps the limits of applying electrostatic physics to climate, showing that heat-accumulation "
   "gradients alone do not determine where catastrophes strike.", 19, False, WHITE)],
], space=1.05)

out = os.path.join(REPO, "Climate_Capacitor_Poster.pptx")
prs.save(out)
print("saved:", out)
